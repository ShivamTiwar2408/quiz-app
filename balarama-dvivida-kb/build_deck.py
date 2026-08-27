#!/usr/bin/env python3
"""Build deck.pptx for the SB 10.67 guide. Palette matches guide/index.html."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG      = RGBColor(0x0B, 0x0D, 0x14)
PANEL   = RGBColor(0x14, 0x1A, 0x2B)
PANEL2  = RGBColor(0x1B, 0x23, 0x38)
LINE    = RGBColor(0x24, 0x2D, 0x45)
INK     = RGBColor(0xEE, 0xF1, 0xF7)
MUTED   = RGBColor(0xA6, 0xB1, 0xC8)
FAINT   = RGBColor(0x6D, 0x77, 0x91)
GOLD    = RGBColor(0xD9, 0xA4, 0x41)   # varuni
RED     = RGBColor(0xC2, 0x62, 0x4A)   # gairika
GREEN   = RGBColor(0x5F, 0x9C, 0x78)   # sala
PURPLE  = RGBColor(0x8E, 0x86, 0xBB)

SERIF = "Georgia"
SANS  = "Calibri"
MONO  = "Consolas"

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def txt(slide, l, t, w, h, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.paragraphs[0].alignment = align
    return tf


def line(tf, text, size, color, *, font=SANS, bold=False, italic=False,
         space_before=0, space_after=0, first=False, align=None, spacing=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if align is not None:
        p.alignment = align
    if spacing:
        p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = font
    r.font.bold = bold
    r.font.italic = italic
    return p


def box(slide, l, t, w, h, fill, border=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, bw=1.2):
    s = slide.shapes.add_shape(shape, l, t, w, h)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if border is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = border
        s.line.width = Pt(bw)
    s.text_frame.word_wrap = True
    return s


def new(eyebrow=None, title=None, kicker=None):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    # gold hairline down the left gutter — the deck's one recurring device
    r = box(s, Inches(0.62), Inches(0.55), Pt(2.2), Inches(6.4), GOLD, shape=MSO_SHAPE.RECTANGLE)
    r.fill.fore_color.rgb = GOLD
    if eyebrow:
        tf = txt(s, Inches(0.95), Inches(0.52), Inches(11.6), Inches(0.32))
        line(tf, eyebrow, 13, GOLD, font=SERIF, italic=True, first=True)
    if title:
        tf = txt(s, Inches(0.95), Inches(0.88), Inches(11.6), Inches(1.0))
        line(tf, title, 33, INK, font=SERIF, first=True, spacing=0.95)
    if kicker:
        tf = txt(s, Inches(0.95), Inches(1.86), Inches(11.4), Inches(0.5))
        line(tf, kicker, 15, MUTED, font=SANS, italic=True, first=True, spacing=1.15)
    return s


def bullets(slide, items, top=2.45, left=0.95, width=11.5, size=16, gap=9):
    tf = txt(slide, Inches(left), Inches(top), Inches(width), Inches(4.4))
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            head, body = it
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(0 if i == 0 else gap)
            p.line_spacing = 1.18
            r = p.add_run(); r.text = "— " + head + "  "
            r.font.size = Pt(size); r.font.color.rgb = INK; r.font.name = SANS; r.font.bold = True
            r2 = p.add_run(); r2.text = body
            r2.font.size = Pt(size); r2.font.color.rgb = MUTED; r2.font.name = SANS
        else:
            line(tf, "— " + it, size, MUTED, first=(i == 0),
                 space_before=(0 if i == 0 else gap), spacing=1.18)
    return tf


def callout(slide, label, text, color, top, left=0.95, width=11.5, height=1.0):
    box(slide, Inches(left), Inches(top), Inches(width), Inches(height), PANEL, LINE)
    bar = box(slide, Inches(left), Inches(top), Pt(3), Inches(height), color, shape=MSO_SHAPE.RECTANGLE)
    bar.fill.fore_color.rgb = color
    tf = txt(slide, Inches(left + 0.22), Inches(top + 0.16), Inches(width - 0.45), Inches(height - 0.3))
    p = tf.paragraphs[0]; p.line_spacing = 1.15
    r = p.add_run(); r.text = label + "  "
    r.font.size = Pt(14); r.font.color.rgb = color; r.font.name = SANS; r.font.bold = True
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(14); r2.font.color.rgb = MUTED; r2.font.name = SANS


def verse(slide, text, cite, top, left=0.95, width=11.5, height=1.35):
    box(slide, Inches(left), Inches(top), Inches(width), Inches(height), PANEL2, LINE)
    tf = txt(slide, Inches(left + 0.28), Inches(top + 0.18), Inches(width - 0.55), Inches(height - 0.3))
    line(tf, text, 16, RGBColor(0xDB, 0xE2, 0xEE), font=SERIF, italic=True, first=True, spacing=1.15)
    line(tf, cite.upper(), 10, FAINT, font=MONO, space_before=8)


def table(slide, rows, top, left=0.95, width=11.5, col_w=None, size=13, row_h=0.44):
    nr, nc = len(rows), len(rows[0])
    shp = slide.shapes.add_table(nr, nc, Inches(left), Inches(top), Inches(width), Inches(row_h * nr))
    tb = shp.table
    tb.first_row = True
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tb.columns[i].width = Emu(int(Inches(width) * cw / total))
    for ri, row in enumerate(rows):
        tb.rows[ri].height = Inches(row_h)
        for ci, cell_text in enumerate(row):
            c = tb.cell(ri, ci)
            c.fill.solid()
            c.fill.fore_color.rgb = PANEL2 if ri == 0 else (PANEL if ri % 2 else BG)
            c.margin_left = c.margin_right = Inches(0.11)
            c.margin_top = c.margin_bottom = Inches(0.05)
            c.vertical_anchor = MSO_ANCHOR.TOP
            tf = c.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.line_spacing = 1.08
            r = p.add_run(); r.text = cell_text
            r.font.size = Pt(10.5 if ri == 0 else size)
            r.font.name = MONO if ri == 0 else SANS
            r.font.bold = ri == 0
            r.font.color.rgb = FAINT if ri == 0 else (INK if ci == 0 else MUTED)
    return tb


# ─────────────────────────────────────────── 1 · title
s = prs.slides.add_slide(BLANK)
s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
box(s, Inches(0), Inches(0), W, Pt(3), GOLD, shape=MSO_SHAPE.RECTANGLE).fill.fore_color.rgb = GOLD
tf = txt(s, Inches(1.1), Inches(1.5), Inches(11.2), Inches(0.4))
line(tf, "CANTO 10  ·  CHAPTER 67  ·  ŚRĪMAD-BHĀGAVATAM", 13, GOLD, font=MONO, first=True)
tf = txt(s, Inches(1.1), Inches(2.15), Inches(11.2), Inches(1.9))
line(tf, "A chapter about a monster", 50, INK, font=SERIF, first=True, spacing=0.95)
line(tf, "is really a chapter about you", 50, GOLD, font=SERIF, italic=True, spacing=0.95)
tf = txt(s, Inches(1.1), Inches(4.35), Inches(11.2), Inches(0.5))
line(tf, "Śrīmad-Bhāgavatam 10.67 — Lord Balarāma slays Dvivida", 21, MUTED, font=SERIF, italic=True, first=True)
box(s, Inches(1.1), Inches(5.05), Inches(4.2), Pt(1), LINE, shape=MSO_SHAPE.RECTANGLE).fill.fore_color.rgb = LINE
tf = txt(s, Inches(1.1), Inches(5.35), Inches(11.2), Inches(1.2))
line(tf, "Fort Lauderdale, FL  ·  Balarāma-pūrṇimā  ·  23 August 2024", 14, MUTED, first=True)
line(tf, "A study guide in 21 slides", 13, FAINT, font=MONO, space_before=7)

# ─────────────────────────────────────────── 2 · the one idea
s = new("sāra", "The one idea everything hangs off",
        "On the surface: thirty verses about a huge ape wrecking a coastline until Balarāma kills him. But the commentators hardly write about the fight. They write about one question — how did a servant of Lord Rāmacandra become a demon?")
verse(s, "Dvivida did not fall because he was weak. He fell because he was strong, and proud of it, and that pride made him slight the Lord's representative.",
      "the one idea", 3.35, height=1.0)
bullets(s, [
    ("Offending a devotee", "takes away your spiritual strength."),
    ("Without that strength,", "you cannot resist bad association."),
    ("Bad association", "gives you anarthas — unwanted things, in the way."),
    ("Anarthas", "bring more offences. This loop has no bottom."),
], top=4.6, size=15, gap=8)
callout(s, "So what is the chapter about?", "The fight is only the setting. That chain is the chapter.", RED, 6.25, height=0.7)

# ─────────────────────────────────────────── 3 · where in the text
s = new("krama", "Where we are in the text",
        "Every chapter of the Bhāgavatam answers a question somebody actually asked. Here the person asking is Mahārāja Parīkṣit. He has seven days left to live, so he does not waste words.")
verse(s, "“I wish to hear further about Śrī Balarāma, the unlimited and immeasurable Supreme Lord, whose activities are all astounding. What else did He do?”",
      "Parīkṣit's request · SB 10.67.1", 2.6, height=1.05)
table(s, [
    ["CHAPTER", "WHAT HAPPENS", "WHY IT MATTERS HERE"],
    ["10.65", "Balarāma visits Vṛndāvana after a long time away. He dances with the gopīs, drinks vāruṇī, and drags the Yamunā off her course with His plough.", "This is the same scene chapter 67 repeats. Dvivida walks into a scene we have already been shown."],
    ["10.66", "Kṛṣṇa kills Pauṇḍraka, the impostor who dressed as Viṣṇu and told Kṛṣṇa to hand over the insignia.", "Balarāma is still in Vṛndāvana, so Parīkṣit hears nothing about Him — which is why he asks again in 67.1."],
    ["10.59", "Eight chapters earlier: Kṛṣṇa kills Narakāsura and frees 16,100 imprisoned princesses.", "This is the death that gives Dvivida his motive. He is not a random monster. He is a friend taking revenge for a friend."],
], top=3.95, col_w=[1.1, 4.4, 5.0], row_h=0.72)

# ─────────────────────────────────────────── 4 · where in time
s = new("kāla", "Where we are in time — and why it matters",
        "Dvivida serves Sugrīva in Rāmacandra's pastimes and dies during Kṛṣṇa's. So if those two appearances are far apart, the chapter's first fact about him turns out to be a fact about his lifespan.")
# the clock, drawn
y = Inches(2.7)
box(s, Inches(0.95), y, Inches(11.5), Inches(0.5), PANEL, LINE)
for i in range(1, 14):
    x = Inches(0.95) + Emu(int(Inches(11.5) * i / 14))
    box(s, x, y, Pt(0.8), Inches(0.5), LINE, shape=MSO_SHAPE.RECTANGLE).fill.fore_color.rgb = LINE
seg = box(s, Inches(0.95) + Emu(int(Inches(11.5) * 6 / 14)), y, Emu(int(Inches(11.5) / 14)),
          Inches(0.5), PANEL2, INK, bw=1.6)
tf = txt(s, Inches(0.95), Inches(2.42), Inches(11.5), Inches(0.26))
line(tf, "ONE DAY OF BRAHMĀ (KALPA) — 1,000 FOUR-YUGA CYCLES · 14 MANUS · ≈71 EACH", 10.5, FAINT, font=MONO, first=True)
tf = txt(s, Inches(0.95) + Emu(int(Inches(11.5) * 5.1 / 14)), Inches(3.28), Inches(3.2), Inches(0.3))
line(tf, "Manu 7 — Vaivasvata", 12, INK, font=SANS, bold=True, first=True)

y2 = Inches(4.05)
box(s, Inches(0.95), y2, Inches(11.5), Inches(0.5), PANEL, LINE)
tf = txt(s, Inches(0.95), Inches(3.78), Inches(11.5), Inches(0.26))
line(tf, "71 FOUR-YUGA CYCLES OF VAIVASVATA MANU", 10.5, FAINT, font=MONO, first=True)
for cyc, col, lbl in ((24, GOLD, "Rāmacandra · cycle 24 · Tretā"), (28, INK, "Kṛṣṇa & Balarāma · cycle 28 · end of Dvāpara")):
    x = Inches(0.95) + Emu(int(Inches(11.5) * (cyc - .5) / 71))
    box(s, x, y2, Emu(int(Inches(11.5) / 71)), Inches(0.5), col, shape=MSO_SHAPE.RECTANGLE).fill.fore_color.rgb = col
span = box(s, Inches(0.95) + Emu(int(Inches(11.5) * 23.5 / 71)), Inches(4.62),
           Emu(int(Inches(11.5) * 4.5 / 71)), Pt(6), RED, shape=MSO_SHAPE.RECTANGLE)
span.fill.fore_color.rgb = RED
tf = txt(s, Inches(4.7), Inches(4.58), Inches(7.7), Inches(0.6))
line(tf, "Dvivida is alive through all of this — about 4 yuga-cycles, tens of millions of years.", 14, RED, first=True)
line(tf, "Rāmacandra: cycle 24, Tretā-yuga.    Kṛṣṇa & Balarāma: cycle 28, end of Dvāpara.", 12, FAINT, font=MONO, space_before=4)
callout(s, "Cirañjīva —", "a being granted exceptional longevity, persisting across yuga-cycles. So is Jāmbavān, who fights in Rāma's pastimes and again in Kṛṣṇa's, and whose daughter Jāmbavatī becomes one of Kṛṣṇa's principal queens. The same lifespan made one the Lord's father-in-law and the other His enemy — longevity is not the deciding variable.", INK, 5.55, height=1.35)

# ─────────────────────────────────────────── 5 · born to serve
s = new("janma", "Born to serve Rāma",
        "Dvivida was not born a demon. He was born on purpose, to help the Lord.")
bullets(s, [
    ("The loophole.", "Rāvaṇa had asked Brahmā to make him unkillable by demigods, demons, snakes and the rest — and forgot to mention animals. So Viṣṇu told the demigods: have sons in the shapes he left out."),
    ("That is why the demigods fathered monkeys.", "The Aśvinī-kumāras had two sons, Mainda and Dvivida. Brothers. Both became ministers to Sugrīva, the monkey king."),
    ("They helped look for Sītā.", "Both went south with Aṅgada's party of thirty thousand monkeys, Hanumān among them — the party that Sampāti, Jaṭāyu's elder brother, finally sent to the Aśoka grove in Laṅkā."),
    ("And he was one of the five asked to stay.", "Before Rāmacandra walked into the Sarayū river, He named five devotees who were to remain: Hanumān, Jāmbavān, Vibhīṣaṇa, Mainda — and Dvivida."),
], top=2.55, size=15, gap=11)
callout(s, "This is the fact the whole chapter turns on.", "The Lord Himself picked him to stay. So everything that happens next happens to someone the Lord chose — which is why “could this happen to me?” is a real question, not a rhetorical one.", RED, 5.75, height=1.15)

# ─────────────────────────────────────────── 6 · strength
s = new("bala", "The measure of his strength",
        "The Bhāgavatam does not just say he was strong. It gives numbers, twice. That strength is what his pride was made of, and the pride is the first link in the chain.")
tf = txt(s, Inches(0.95), Inches(2.55), Inches(11.5), Inches(0.3))
line(tf, "THE JUMPING CONTEST AT THE SOUTHERN SHORE — WHO CAN CROSS 100 YOJANAS (≈800 MILES)?", 10.5, FAINT, font=MONO, first=True)
bars = [("Gaja", 10, "10 yojanas · 80 mi", RGBColor(0x3D, 0x4A, 0x6B)),
        ("Gavākṣa", 20, "20 yojanas · 160 mi", RGBColor(0x3D, 0x4A, 0x6B)),
        ("Mainda", 60, "60 yojanas · 480 mi — “I am capable”", PURPLE),
        ("Dvivida", 70, "70 yojanas · 560 mi — “Mainda was not as strong as me”", RED),
        ("Hanumān", 100, "“a hundred yojanas — and back again.” The only one who did not compare himself.", INK)]
bx, bw_full = Inches(2.35), Inches(6.4)
for i, (name, val, lbl, col) in enumerate(bars):
    t = Inches(2.92 + i * 0.44)
    tfn = txt(s, Inches(0.95), t + Inches(0.03), Inches(1.3), Inches(0.3), align=PP_ALIGN.RIGHT)
    line(tfn, name, 13, INK, bold=True, first=True)
    b = box(s, bx, t, Emu(int(bw_full * val / 100)), Inches(0.26), col, shape=MSO_SHAPE.RECTANGLE)
    b.fill.fore_color.rgb = col
    tfl = txt(s, bx + Emu(int(bw_full * val / 100)) + Inches(0.12), t + Inches(0.02), Inches(4.6), Inches(0.3))
    line(tfl, lbl, 11.5, col if col not in (RGBColor(0x3D, 0x4A, 0x6B),) else FAINT, font=MONO, first=True)
callout(s, "Listen to how they say it.", "Mainda says what he can do. Dvivida says he can do more than Mainda. That comparison is the moment strength turns into pride in strength — bala into bala-abhimāna.", GOLD, 5.3, height=1.05)
callout(s, "The second number —", "by this chapter he churns the ocean with the strength of ten thousand elephants, pulls up mountains, and rips down śāla hardwood with one hand. SB 10.2 already lists him among the demons Kaṁsa could call on.", MUTED, 6.45, height=0.85)

# ─────────────────────────────────────────── 7 · eternal associate
s = new("nitya-siddha", "Was he an eternal associate? The hard question",
        "The commentaries say Mainda and Dvivida are eternally liberated — attendants worshipped along with Rāmacandra. But souls in the Lord's abode don't fall down. So what happened here?")
bullets(s, [
    ("Jīva Gosvāmī, Bhakti-sandarbha:", "they are worshipped as āvaraṇa-devatās — attendants who surround the Lord and belong to His entourage, without showing Viṣṇu's own fullness."),
    ("Jīva Gosvāmī on this very chapter:", "“…persons endowed with śakti from the spiritual realm. Because they did not respect Lakṣmaṇa — like Jaya and Vijaya — they appeared as demons in order to show the result of committing offences.”"),
    ("The distinction that resolves it:", "some jīvas are given a portion of an eternal associate's power and carry the same name. Those jīvas can be corrupted, because they are still jīvas."),
    ("Viśvanātha Cakravartī Ṭhākura:", "“the Lord arranged their degradation to show the evil of the bad association that results from offending great personalities.”"),
], top=2.6, size=14.5, gap=11)
callout(s, "Please don't turn this into a general theory.", "We were not all once living in Vaikuṇṭha and then fell. Jaya and Vijaya were an exception arranged to teach us something, and so is Dvivida — “it's the exception, it's not the rule.” Take the lesson, not a theory about where you came from.", RED, 6.05, height=1.15)

# ─────────────────────────────────────────── 8 · the two offences
s = new("aparādha", "The two offences",
        "Pride by itself was not the offence. Pride was the condition he was in. The offences are what the pride made him say and think — two of them, both about Lakṣmaṇa, both easy to read past in the Rāmāyaṇa.")
tf = txt(s, Inches(0.95), Inches(2.6), Inches(5.5), Inches(0.3))
line(tf, "OFFENCE ONE — THE INTRODUCTION", 11, RED, font=MONO, first=True)
box(s, Inches(0.95), Inches(2.95), Inches(5.5), Inches(1.55), PANEL, RED)
tf = txt(s, Inches(1.15), Inches(3.12), Inches(5.1), Inches(1.25))
line(tf, "“Give your full, undivided attention to Rāma. The other fellow with Him — Lakṣmaṇa — is not so important.”", 15, RGBColor(0xDB, 0xE2, 0xEE), font=SERIF, italic=True, first=True, spacing=1.12)
line(tf, "Said to Sugrīva while introducing him to Rāma — while doing service, and meaning well.", 12, FAINT, space_before=8)

tf = txt(s, Inches(6.9), Inches(2.6), Inches(5.5), Inches(0.3))
line(tf, "OFFENCE TWO — THE BATTLEFIELD JUDGMENT", 11, RED, font=MONO, first=True)
box(s, Inches(6.9), Inches(2.95), Inches(5.5), Inches(1.55), PANEL, RED)
tf = txt(s, Inches(7.1), Inches(3.12), Inches(5.1), Inches(1.25))
line(tf, "“What kind of great hero is this?”", 15, RGBColor(0xDB, 0xE2, 0xEE), font=SERIF, italic=True, first=True, spacing=1.12)
line(tf, "Thought, not said. Lakṣmaṇa was lying in great pain, hit by the śakti weapon. Nobody heard this.", 12, FAINT, space_before=8)

callout(s, "Neither one was hostile. Neither one was heard.", "The first was friendly advice. The second was a private thought. Both still counted. Offending a devotee does not need anger, volume, or an audience.", RED, 4.85, height=1.0)
callout(s, "Both come from the same place.", "He measured Lakṣmaṇa against himself — first not important enough, then not heroic enough. That is what pride actually does: it makes you the yardstick.", GOLD, 6.0, height=1.0)

# ─────────────────────────────────────────── 9 · why Laksmana
s = new("guru-tattva", "Why Lakṣmaṇa, of all people",
        "If we shorten this to “don't offend Lakṣmaṇa,” it is easy to answer: fine, I'm not likely to meet him. But Lakṣmaṇa is here as an example of something much wider.")
verse(s, "Guru-tattva — the Lord's representative. The one who introduces you to the Supreme Lord, instructs you, and gives you His association. Lakṣmaṇa is that to Rāma. Balarāma is that to Kṛṣṇa. So is Nanda Mahārāja.",
      "the category, not the person", 2.7, height=1.3)
bullets(s, [
    ("What he actually did:", "he pointed everyone towards the Lord and put down the person who was making that introduction possible. He kept the Lord and threw away the connection to Him."),
    ("And look how it comes back:", "he offends the guru-tattva in Rāma's pastimes, and is killed by the guru-tattva in Kṛṣṇa's. Yugas later, the same principle he brushed aside is the one that deals with him."),
], top=4.3, size=16, gap=12)
callout(s, "The point I want to insist on:", "“It's not just that you're not Lakṣmaṇa, so I can mess with you.” The principle is vaiṣṇava-aparādha — offending any devotee — and you can lose your spiritual standing over it.", RED, 5.95, height=1.1)

# ─────────────────────────────────────────── 10 · the mechanism
s = new("anartha", "The mechanism of degradation",
        "Viśvanātha Cakravartī Ṭhākura gives the whole sequence in one sentence, and every step causes the next. Watch the order carefully: bad association is the punishment for the offence, not the reason for it.")
steps = [("Attendant in the Lord's abode", "āvaraṇa-devatā · eternally situated", INK),
         ("Minister of Sugrīva, servant of Rāma", "made for this service", INK),
         ("Pride in his own strength", "bala-abhimāna · the condition, not yet an offence", GOLD),
         ("Offence I — “Lakṣmaṇa is not important”", "spoken while rendering service", RED),
         ("Offence II — “What kind of hero is this?”", "thought only in the mind", RED),
         ("His spiritual strength is taken away", "the reaction · this is literal", RED),
         ("Now bad association can get at him", "→ close friendship with Narakāsura", RED),
         ("Unwanted things fill his heart", "revenge · arson · sacrilege · molesting women", RED),
         ("More offences, with no bottom to it", "and round it goes again", RED)]
for i, (head, sub, col) in enumerate(steps):
    t = Inches(2.5 + i * 0.47)
    l = Inches(0.95 + i * 0.30)
    b = box(s, l, t, Inches(5.0), Inches(0.42), PANEL, col, bw=1.3)
    tf = txt(s, l + Inches(0.16), t + Inches(0.05), Inches(4.7), Inches(0.34))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = head + "   "
    r.font.size = Pt(12.5); r.font.color.rgb = INK; r.font.name = SANS; r.font.bold = True
    r2 = p.add_run(); r2.text = sub
    r2.font.size = Pt(10); r2.font.color.rgb = FAINT; r2.font.name = MONO
box(s, Inches(8.4), Inches(2.5), Inches(4.0), Inches(1.35), PANEL2, INK, bw=1.8)
tf = txt(s, Inches(8.62), Inches(2.68), Inches(3.6), Inches(1.05))
line(tf, "Back home — brahmajyoti or Vaikuṇṭha", 14, INK, bold=True, first=True, spacing=1.1)
line(tf, "Killed by Balarāma's bare hands, and liberated by it. SB 2.7.34–35.", 12, MUTED, space_before=6, spacing=1.1)
tf = txt(s, Inches(8.4), Inches(4.2), Inches(4.0), Inches(2.4))
line(tf, "He went all the way down, and he still got back.", 15, GOLD, font=SERIF, italic=True, first=True, spacing=1.15)
line(tf, "There is only one step on this chain that you choose — the offence. After that, it runs on its own.", 13, MUTED, space_before=12, spacing=1.2)
line(tf, "Once your spiritual strength is gone, resisting bad association is not a question of willpower. Willpower is the thing that was taken away.", 13, MUTED, space_before=9, spacing=1.2)
line(tf, "Only the first three steps are ones you can do anything about.", 13, RED, bold=True, space_before=9, spacing=1.2)

# ─────────────────────────────────────────── 11 · Narakasura
s = new("saṅga", "Narakāsura: the association",
        "The friend who ruined Dvivida had been ruined the same way himself.")
bullets(s, [
    ("He was not born a demon either.", "His name was Naraka. His mother was Bhūmi — the earth herself. His father was Lord Varāha. You cannot get better parents than that."),
    ("He went bad the same way Dvivida did —", "by bad association. By SB 10.59 he is Narakāsura, keeping 16,100 princesses locked up. They appealed to Kṛṣṇa, and to free them Kṛṣṇa had to kill him."),
    ("Then Dvivida takes his decision.", "“He killed my best friend. I am going to get back at him.” And he set out to wreck Kṛṣṇa's flourishing kingdom."),
], top=2.6, size=15.5, gap=11)
callout(s, "Notice how it passes along.", "Naraka was spoiled by bad association. Then Dvivida was spoiled by associating with Naraka. It spreads — which is why association is treated as a serious spiritual matter and not just a social one.", GOLD, 4.95, height=1.05)
callout(s, "And notice why he does it.", "This is not hunger and it is not madness. It is revenge, for somebody he loved. That is the hardest kind to catch in yourself, because it never feels like a fault. It feels like loyalty.", GREEN, 6.1, height=1.0)

# ─────────────────────────────────────────── 12 · the rampage
s = new("Ānarta", "The rampage on Ānarta",
        "He went for the region around Dvārakā. The Bhāgavatam lists the damage one item at a time, and it is worth reading slowly — this is what a chain of unwanted things looks like once nobody is stopping it.")
bullets(s, [
    "He set fires. Cities burned, villages burned, mines burned, the cowherds' settlements burned.",
    "He pulled up whole mountains and threw them at the neighbouring kingdoms — Ānarta most of all, because that is where Hari lives.",
    "He walked into the ocean and churned the water with his arms until the coast went under.",
    "He tore down the trees around the sages' huts, and fouled their sacrificial fires with his excrement and urine. He wore a belt of skulls.",
    "He shut men and women into mountain caves and rolled boulders across the mouths — “as a wasp imprisons smaller insects.”",
    "He took an enormous form, looted the crops until people went hungry, and polluted women from respectable families.",
], top=2.6, size=14.5, gap=7)
table(s, [
    ["PLACE", "WHAT IT IS"],
    ["Dvārakā", "The island capital where Kṛṣṇa and Balarāma reside; Ugrasena rules over subordinate kings."],
    ["Prabhāsa", "A holy place on the coast, within the same province."],
    ["Raivataka Mountain", "In King Raivata's province — where Balarāma is found dancing. Revatī, Balarāma's wife, is his daughter."],
    ["Ānarta vs. anartha", "Ānarta is the province. Anartha is the unwanted thing. Prabhupāda once named a disciple Ānartā, to general alarm — until the province was recalled."],
], top=4.9, col_w=[2.4, 9.1], row_h=0.48, size=12)

# ─────────────────────────────────────────── 13 · Raivataka
s = new("līlā", "Raivataka: the scene he walked into",
        "He heard sweet singing coming from Raivataka Mountain, so he went to see what it was. What he found was Balarāma enjoying Himself — and every detail of that scene is about to become a target.")
verse(s, "There he saw Śrī Balarāma, adorned with a garland of lotuses and appearing most attractive in every limb. He was singing amidst a crowd of young women, and since He had drunk vāruṇī liquor, His eyes rolled as if He were intoxicated. His body shone brilliantly as He behaved like an elephant in mada.",
      "the scene at Raivataka · SB 10.67", 2.65, height=1.45)
table(s, [
    ["ELEMENT", "WHAT IT IS", "WHY IT'S THERE"],
    ["Revatī", "Daughter of King Raivata (Kakudmī), and Balarāma's wife. One of the women there.", "That marriage is the reason Balarāma is on this mountain in the first place."],
    ["Vāruṇī", "A fragrant drink. Two chapters earlier Varuṇa had it collect in a tree hollow so Balarāma could dance.", "How it turns up again here is not explained. It just does."],
    ["Mada", "The scented fluid that runs from a bull elephant's temples when he is in rut.", "The verse's own comparison for great strength in a happy, unguarded mood."],
    ["Plough & club Sunanda", "Both there — but not in His hands. He is dancing.", "That is exactly what makes the next scene possible."],
], top=4.25, col_w=[2.3, 5.2, 4.0], row_h=0.6, size=12)

# ─────────────────────────────────────────── 14 · the provocation
s = new("avajñā", "The provocation",
        "Watch what he does, because he does not attack. He heckles — and each thing he does is one step worse than the last.")
bullets(s, [
    "He climbs onto a branch, shakes the trees, and makes the sound kilakilā — a gorilla sound.",
    "The girls with Balarāma laugh at his cheek. They are young, they like a joke, they are a little silly.",
    "Being laughed at is what sets him off. Now he insults the girls — pulling faces with his eyebrows, walking about absurdly, exposing himself and urinating, and all of this while Balarāma is watching him.",
    "Viśvanātha Cakravartī Ṭhākura points out the worst part: he did not even glance at Balarāma.",
    "Balarāma throws a rock at him. He dodges it and grabs the pot of vāruṇī.",
    "He smashes the pot, then starts pulling at the girls' clothes, laughing at Him the whole time.",
], top=2.5, size=14, gap=7)
callout(s, "The one detail worth this whole slide.", "Balarāma's first move is a rock — not His club. The commentary says why: because of Dvivida's previous position as a devotee. Even now, He remembers who this used to be.", RED, 4.75, height=1.05)
callout(s, "And this is not what He decides on.", "Only after seeing all this rudeness does Balarāma think of the wrecked kingdoms all around — and then He picks up His club and plough. He decides on the ruined province, not the ruined party.", GOLD, 5.9, height=1.05)

# ─────────────────────────────────────────── 15 · the fight
s = new("yuddha", "The fight, blow by blow",
        "Dvivida keeps reaching for something bigger until he runs out of things to throw, while Balarāma keeps putting weapons down until He has none left.")
tf = txt(s, Inches(0.95), Inches(2.5), Inches(5.6), Inches(0.3))
line(tf, "DVIVIDA — ESCALATING", 11, RED, font=MONO, first=True)
tf = txt(s, Inches(6.85), Inches(2.5), Inches(5.6), Inches(0.3))
line(tf, "BALARĀMA — DISARMING", 11, INK, font=MONO, first=True)
rounds = [("Pulls up a śāla tree with one hand and hits Him on the head", "Does not move — stands like a mountain — and catches the tree"),
          ("Does not even notice his own wound. Another tree, then another", "Hits his skull with the club Sunanda — “a mountain beautified by streaks of red oxide”"),
          ("No trees left in the forest, so he throws stones", "Breaks up every tree he throws, then every stone"),
          ("Closes his fists and beats on the Lord's body", "Puts club and plough aside — bare hands only, arms like palm trees"),
          ("Falls down vomiting blood. The mountain shakes “like a boat tossed about at sea.”", "One blow to the collarbone")]
for i, (l_txt, r_txt) in enumerate(rounds):
    t = Inches(2.88 + i * 0.72)
    last = i == len(rounds) - 1
    box(s, Inches(0.95), t, Inches(5.6), Inches(0.62), PANEL, RED, bw=1.6 if last else 1.1)
    tf = txt(s, Inches(1.12), t + Inches(0.1), Inches(5.3), Inches(0.45))
    line(tf, l_txt, 13, INK if last else MUTED, bold=last, first=True, spacing=1.08)
    box(s, Inches(6.85), t, Inches(5.6), Inches(0.62), PANEL2 if last else PANEL, INK, bw=1.8 if last else 1.1)
    tf = txt(s, Inches(7.02), t + Inches(0.1), Inches(5.3), Inches(0.45))
    line(tf, r_txt, 13, INK if last else MUTED, bold=last, first=True, spacing=1.08)
tf = txt(s, Inches(0.95), Inches(6.62), Inches(11.5), Inches(0.6))
line(tf, "The two sides move in opposite directions. He goes tree → stone → fist; the Lord goes club → nothing at all — and that is when it finishes. The demigods called out, “Victory to You! Well done!” and showered flowers.", 13.5, GOLD, font=SERIF, italic=True, first=True)
line(tf, "Then He let out the people who had been sealed in the caves, and went back to His capital while everyone sang His glories.", 12, FAINT, space_before=5)

# ─────────────────────────────────────────── 16 · liberation
s = new("mukti", "Killed, and therefore liberated",
        "There is one point at the end that changes how you read every killing in the Tenth Canto. It comes from SB 2.7, where the Lord's coming appearances are listed in advance.")
verse(s, "All these demoniac personalities — Pralamba, Dhenuka, Baka, Keśī, Ariṣṭa, Cāṇūra, Muṣṭika, Kuvalayāpīḍa, Kaṁsa, Yavana, Narakāsura, Pauṇḍraka; marshals like Śālva, Dvivida the monkey, Balvala, Dantavakra, the seven bulls, Śambara, Vidūratha and Rukmī — would all fight vigorously with the Lord Hari, or with Him under the names Baladeva, Arjuna, Bhīma; and the demons, being thus killed, would attain either the impersonal brahmajyoti or His personal abode in the Vaikuṇṭha planets.",
      "SB 2.7.34–35 · the scheduled incarnations", 2.6, height=1.9)
callout(s, "So he went down, and he came back.", "The soul who lost his place as an attendant in the Lord's abode is put back there — by being killed by the Lord.", INK, 4.75, height=0.8)
callout(s, "“We don't imitate.”", "Let me stop the obvious conclusion right here. You cannot arrange to be the exception — that is what exception means. The mercy Dvivida got tells us something about the Lord. It is not a plan any of us can follow.", RED, 5.7, height=1.0)
callout(s, "Why He comes at all — BG 4.7–8.", "Whenever religion declines and irreligion rises, He comes: to protect the devotees, finish off the miscreants, and re-establish religion, age after age. Balarāma frees the villagers from the caves and kills the ape on the same afternoon.", GOLD, 6.8, height=0.62)

# ─────────────────────────────────────────── 17 · kavaca
s = new("kavaca", "Balarāma's four protections",
        "There is a prayer for protection in the Garga-saṁhitā that takes each demon Balarāma killed and pairs it with an enemy inside us. The pastime becomes a prayer, and the prayer tells you what is wrong with you.")
table(s, [
    ["MAY LORD BALARĀMA…", "…PROTECT ME FROM", "WHY THESE GO TOGETHER"],
    ["the enemy of Dhenukāsura", "lust · kāma", "Lust burns like a fire. Once it is lit, whatever you feed it makes it bigger."],
    ["who killed Dvivida", "anger · krodha", "Our chapter. Revenge drives Dvivida from beginning to end — and once anger has you, you cannot stop it from inside."],
    ["the enemy of Balvala", "greed · lobha", "Balvala fouled the sacrifices at Naimiṣāraṇya — the appetite that spoils the thing it feeds on."],
    ["the enemy of Jarāsandha", "illusion · moha", "Jarāsandha attacked again and again, believing every time that this time it would work."],
], top=2.75, col_w=[3.0, 2.3, 6.2], row_h=0.66)
callout(s, "Here is what the prayer is doing:", "Balarāma's victories are not just old news. Each one shows something He can do, and the prayer asks Him to do that same thing inside you. He has beaten anger — in public, on Raivataka Mountain — so He is the one to ask about yours.", INK, 6.15, height=1.05)

# ─────────────────────────────────────────── 18 · the tongue
s = new("vāk-saṁyama", "Controlling the tongue",
        "A fair question to ask in a chapter where the first offence was a sentence somebody said. Here are four things that help.")
bullets(s, [
    ("1 · Build up the mode of goodness.", "BG 17.15 says speech should be truthful and pleasing. Both halves are the instruction. And goodness is built by habits — cleanliness, compassion, charity, truthfulness — not switched on when you need it."),
    ("2 · The tongue does two jobs —", "it tastes and it makes sound. So eat only kṛṣṇa-prasādam. Not as a rule to obey: prasādam purifies the tongue, and then control follows from the purity instead of the other way round."),
    ("3 · Learn how to talk to people.", "Non-violent communication, which devotees are taught as empathic communication. The main idea: know your own needs and take steps to meet them, instead of expecting others to. Unmet expectation is what starts the tongue going."),
    ("4 · Chant.", "The best shelter of all — and the only one that uses both of the tongue's jobs at the same time."),
], top=2.5, size=14.5, gap=11)
callout(s, "The trap in obeying only half of BG 17.15:", "“I only have to tell the truth, because I'm in the mode of goodness.” That is not the mode of goodness. Goodness includes the pleasing part. A truth used as a weapon has changed modes without changing a word.", RED, 6.1, height=1.1)

# ─────────────────────────────────────────── 19 · anger
s = new("krodha", "Anger, and why bhakti is the cure",
        "The sharpest question I was asked: “I understand once anger kicks in there is no stopping point — so I think I should stop before it kicks in. But it kicks in.”")
bullets(s, [
    ("What BG 2.62–63 is actually describing:", "anger produces delusion, delusion makes you forget yourself, and then intelligence is lost. So the thing you would use to stop it is the thing it switches off."),
    ("The example that proves it — Bhīma and Duryodhana,", "both of them Balarāma's own students in club-fighting. Their duel went on for days. He travelled there and asked them to stop: Bhīma is stronger, Duryodhana is more skilled, neither of you can win — stop wasting your time."),
    ("They could not stop.", "Not would not. Could not. When anger is running you, you cannot take good advice even from the wisest, dearest person you know. Not even from Balarāma. That is how far willpower gets you."),
    ("So change the mode instead of fighting it.", "Anger is ignorance mixed with passion. SB 1.2 gives the way out — worshipping Viṣṇu raises you to goodness. Give Him your possessions, your intelligence, your energy, your work and what it earns."),
], top=2.5, size=14.5, gap=11)
callout(s, "And then just ask Him.", "Not self-improvement — the kavaca: “May Lord Balarāma, who killed Dvivida, always protect me from anger.” He can take the tendency away. He has done it before, in public, to the angriest and strongest person in the chapter. Bhakti is the cure. And don't get angry at Balarāma.", GOLD, 6.15, height=1.05)

# ─────────────────────────────────────────── 20 · where I take a problem
s = new("āśraya", "Where I take a problem",
        "We are reading this on Balarāma-pūrṇimā, so let me finish where the chapter finishes — not with the ape, but with the two brothers, and with something you can actually do about the anger, the pride and the offences.")
bullets(s, [
    ("Śrīla Prabhupāda was asked on a morning walk", "what to do about an obstacle you cannot get past. Go and stand before the deity — Gaurāṅga, or Kṛṣṇa-Balarāma, they are not different — and say, “Sir, I have a problem. Please help me.” And he will do the needful."),
    ("He added two things.", "Not because the Lord is there to supply your orders — that is not the relationship. But because He is kind to someone who depends on Him the way a child depends on a parent. And: don't make a habit of bothering Him too many times."),
    ("Somebody in the room confirmed it.", "A devotee said they pray exactly like that, and it works — not instantaneously, but it works. When you really need something, you know where to go."),
], top=2.5, size=14.5, gap=11)
callout(s, "That is a good pair to have on your side:", "Kṛṣṇa is very kind and Balarāma is all-powerful.", GOLD, 5.95, height=0.95)

# ─────────────────────────────────────────── 21 · close
s = prs.slides.add_slide(BLANK)
s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
box(s, Inches(0), Inches(7.46), W, Pt(3), GOLD, shape=MSO_SHAPE.RECTANGLE).fill.fore_color.rgb = GOLD
tf = txt(s, Inches(1.1), Inches(1.35), Inches(11.2), Inches(0.35))
line(tf, "WHERE TO INTERVENE", 12, GOLD, font=MONO, first=True)
tf = txt(s, Inches(1.1), Inches(1.85), Inches(10.8), Inches(2.2))
line(tf, "Find the last step you can still do something about.", 34, INK, font=SERIF, first=True, spacing=1.0)
line(tf, "It is not bad association — by the time you are there, most of it is already decided. It is pride: the habit of measuring devotees against yourself.",
     19, MUTED, font=SERIF, italic=True, space_before=14, spacing=1.2)
box(s, Inches(1.1), Inches(4.55), Inches(4.2), Pt(1), LINE, shape=MSO_SHAPE.RECTANGLE).fill.fore_color.rgb = LINE
tf = txt(s, Inches(1.1), Inches(4.85), Inches(11.2), Inches(2.2))
line(tf, "SOURCES", 11, FAINT, font=MONO, first=True)
line(tf, "From the class “SB 10.67 — Lord Balarāma slays Dvivida Gorilla,” Fort Lauderdale FL, Balarāma-pūrṇimā, 23 August 2024 (1 hr 17 min).",
     13, MUTED, space_before=8, spacing=1.2)
line(tf, "Śrīmad-Bhāgavatam 10.67 with 10.2, 10.59, 10.65, 10.66, 2.7.34–35, 1.2 · Bhagavad-gītā 2.62–63, 4.7–8, 17.15 · Vālmīki Rāmāyaṇa (Kiṣkindhā- and Uttara-kāṇḍa) · Viśvanātha Cakravartī Ṭhākura on SB 10.67 · Jīva Gosvāmī, Bhakti-sandarbha, his SB commentary and Gopāla-campū · Garga-saṁhitā · Bhaktivinoda Ṭhākura, Caitanya-śikṣāmṛta and Kṛṣṇa-saṁhitā.",
     11.5, FAINT, space_before=7, spacing=1.2)
line(tf, "Full guide: shivamtiwar2408.github.io/quiz-app/balarama-dvivida-kb/guide/index.html", 12.5, GOLD, font=MONO, space_before=10)

prs.save("/tmp/ytg-balarama-dvivida-kb/deck.pptx")
print("slides:", len(prs.slides._sldIdLst))
